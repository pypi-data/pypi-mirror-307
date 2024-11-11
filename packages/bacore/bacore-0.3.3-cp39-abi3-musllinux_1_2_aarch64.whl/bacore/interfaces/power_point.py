"""Power point interface."""

from pptx import Presentation
from pptx.util import Inches
from typing import Optional

WIDESCREEN_WIDTH = Inches(13.33)
WIDESCREEN_HEIGHT = Inches(7.5)


class PowerPoint:
    """PowerPoint class with pptx.Presentation as `prs` attribute.

    Attributes:
        `prs`: Contains the pptx.Presentation object.

    Methods:
        `add_slide`: Add a slide from default tempalte using template index and having an optional title.
    """

    def __init__(self) -> None:
        self.prs = Presentation()

    def add_slide(self, layout_index: int, title_text: Optional[str] = None):
        """Add a pptx slide with an optional title text.

        Returns:
            slide object
        """
        slide_layout = self.prs.slide_layouts[layout_index]
        slide = self.prs.slides.add_slide(slide_layout)
        if title_text:
            title = slide.shapes.title
            title.text = title_text
        return slide


def default_templates(widescreen_dimensions: Optional[bool] = True) -> None:
    """Create a power point slide presentation using default templates."""
    ppt = PowerPoint()

    if widescreen_dimensions:
        ppt.prs.slide_width = WIDESCREEN_WIDTH
        ppt.prs.slide_height = WIDESCREEN_HEIGHT

    [ppt.add_slide(layout_index) for layout_index in range(11)]

    ppt.prs.save("default_templates.pptx")
