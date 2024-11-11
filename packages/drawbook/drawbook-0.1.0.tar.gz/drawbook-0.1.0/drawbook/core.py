"""
Core functionality for the drawbook library.
"""

from pathlib import Path
from typing import List, Literal
import tempfile
import io
import requests
import warnings
from tqdm import tqdm
from PIL import Image
import huggingface_hub
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from huggingface_hub import InferenceClient

class Book:
    """A class representing a children's book that can be exported to PowerPoint."""
    
    def __init__(
        self,
        title: str = "Untitled Book",
        pages: List[str] = None,
        title_illustration: str | Literal[False] | None = None,
        illustrations: List[str | None | Literal[False]] = None,
        lora: str = "SebastianBodza/Flux_Aquarell_Watercolor_v2",
        author: str | None = None,
        illustration_prompts: List[str | None] = None,
        title_illustration_prompt: str | None = None
    ):
        """
        Initialize a new Book.
        
        Args:
            title: The book's title
            pages: List of strings containing text for each page
            illustrations: List of illustration paths or placeholders
                         (str for path, None for pending, False for no illustration)
            lora: The LoRA model on Hugging Face to use for illustrations
            author: The book's author name
            illustration_prompts: Optional list of custom prompts for page illustrations
            title_illustration_prompt: Optional custom prompt for title illustration
        """
        self.title = title
        self.pages = pages or []
        self.illustrations = illustrations or []
        self.lora = lora
        self.title_illustration = title_illustration
        self.author = author
        self.illustration_prompts = illustration_prompts or []
        self.title_illustration_prompt = title_illustration_prompt
        self.client = InferenceClient()
        
        # Ensure illustrations list matches pages length
        while len(self.illustrations) < len(self.pages):
            self.illustrations.append(None)
        
        # Ensure illustration_prompts list matches pages length
        while len(self.illustration_prompts) < len(self.pages):
            self.illustration_prompts.append(None)
            

    def _get_illustration_prompt(self, text: str) -> str:
        """Get an illustration prompt from the text using Qwen."""
        system_prompt = """You are a helpful assistant that converts children's book text into illustration prompts. 
        Extract a key object along with its description that could be used to illustrate the page. 
        Replace any proper names with more generic versions.
        
        For example:
        If the text is: "Mustafa loves his silver cybertruck. One day, his cybertruck starts to glow, grow, and zoom up into the sky"
        You should return: "A silver cybertruck zooming into the sky"
        
        If the text is: "Up, up, up goes Mustafa in his special cybertruck. He waves bye-bye to his house as it gets tiny down below"
        You should return: "A boy in the sky waving bye"
        """
        
        user_prompt = f"""This is the text of a page in a children's book. From this text, extract a key object along with its description that could be used to illustrate this page. Replace any proper names with more generic versions.

Text: {text}

Return ONLY the illustration description, nothing else."""

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ]

        stream = self.client.chat.completions.create(
            model="Qwen/Qwen2.5-72B-Instruct",
            messages=messages,
            max_tokens=500,
            stream=True
        )

        response = ""
        for chunk in stream:
            if chunk.choices[0].delta.content is not None:
                response += chunk.choices[0].delta.content

        return response.strip()

    def _get_prompt(self, illustration_prompt: str) -> str:
        if self.lora == "SebastianBodza/Flux_Aquarell_Watercolor_v2":
            return f"A AQUACOLTOK watercolor painting with a white background of: {illustration_prompt}"
        else:
            warnings.warn(
                f"The LoRA model '{self.lora}' is not officially supported. "
                "Results may not be as expected.",
                UserWarning
            )
            return f"An illustration of: {illustration_prompt}"

    def export(self, filename: str | Path | None = None) -> None:
        """
        Export the book to a PowerPoint file.
        
        Args:
            filename: Optional path where to save the file. If None, creates in temp directory.
        """
        if filename is None:
            # Create temp file with .pptx extension
            temp_file = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
            output_path = Path(temp_file.name)
            temp_file.close()
        else:
            # Convert to Path object and resolve to absolute path
            output_path = Path(filename).resolve()
            # Ensure parent directories exist
            output_path.parent.mkdir(parents=True, exist_ok=True)

        prs = Presentation()
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Add diagonal striped border at the top
        border = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(0.2), Inches(7.5)
        )
        border.fill.solid()
        border.fill.fore_color.rgb = RGBColor(128, 0, 0)  # Maroon
        
        # Add title illustration first if available
        if isinstance(self.title_illustration, str):
            try:
                slide.shapes.add_picture(
                    self.title_illustration,
                    Inches(2.5), Inches(1.5),
                    Inches(5), Inches(5)
                )
            except Exception as e:
                print(f"Warning: Could not add title illustration: {e}")
        
        # Add title with adjusted positioning and z-order
        title = slide.shapes.title
        title.top = Inches(0)
        title.height = Inches(2.0)
        title.width = Inches(10)
        
        # Add title text
        p1 = title.text_frame.paragraphs[0]
        # Clear any existing text
        p1.clear()
        p1.font.name = "Trebuchet MS"
        p1.alignment = PP_ALIGN.CENTER
        
        # Define common stop words
        stop_words = {'a', 'an', 'and', 'are', 'as', 'at', 'be', 'by', 'for',
                     'from', 'has', 'he', 'in', 'is', 'it', 'its', 'of', 'on',
                     'that', 'the', 'to', 'was', 'were', 'will', 'with'}
        
        # Split title and add each word with appropriate size
        words = self.title.split()
        for i, word in enumerate(words):
            run = p1.add_run()
            run.text = word + (' ' if i < len(words) - 1 else '')
            run.font.name = "Trebuchet MS"
            if word.lower() in stop_words:
                run.font.size = Inches(0.42)  # Smaller size for stop words
            else:
                run.font.size = Inches(0.5)   # Regular size for other words
        
        # Add author with adjusted positioning
        if self.author is not None:
            author_box = slide.shapes.add_textbox(
                Inches(0), Inches(6.5),  # Moved up from bottom
                Inches(10), Inches(0.5)
            )
            author_frame = author_box.text_frame
            author_frame.text = f"Written by {self.author}"
            author_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            author_frame.paragraphs[0].font.name = "Trebuchet MS"
            author_frame.paragraphs[0].font.size = Inches(0.25)
        
        # Add content slides
        content_slide_layout = prs.slide_layouts[5]  # Blank layout
        
        for page_num, (text, illustration) in enumerate(zip(self.pages, self.illustrations)):
            slide = prs.slides.add_slide(content_slide_layout)
            
            if isinstance(illustration, str):
                try:
                    slide.shapes.add_picture(
                        illustration,
                        Inches(2.5), Inches(1.4),
                        Inches(5), Inches(5)
                    )
                except Exception as e:
                    print(f"Warning: Could not add illustration on page {page_num + 1}: {e}")
            
            # Split text into sentences and join with newlines
            sentences = text.replace('. ', '.\n').split('\n')
            
            # Special formatting for first page
            if page_num == 0 and text:
                # Split first character from first sentence
                first_char = sentences[0][0]
                first_sentence_rest = sentences[0][1:]
                
                p = slide.shapes.title.text_frame.paragraphs[0]
                p.line_spacing = 1.5  # Add line spacing
                run = p.add_run()
                run.text = first_char
                run.font.size = Inches(0.3)
                run.font.name = "Trebuchet MS"
                
                run = p.add_run()
                run.text = first_sentence_rest
                run.font.size = Inches(0.25)
                run.font.name = "Trebuchet MS"
                
                # Add remaining sentences as new paragraphs
                for sentence in sentences[1:]:
                    p = slide.shapes.title.text_frame.add_paragraph()
                    p.line_spacing = 1.5  # Add line spacing
                    p.text = sentence
                    p.font.name = "Trebuchet MS"
                    p.font.size = Inches(0.25)
                    p.alignment = PP_ALIGN.CENTER
            else:
                # Add each sentence as a separate paragraph
                first_paragraph = True
                for sentence in sentences:
                    if first_paragraph:
                        p = slide.shapes.title.text_frame.paragraphs[0]
                        first_paragraph = False
                    else:
                        p = slide.shapes.title.text_frame.add_paragraph()
                    p.line_spacing = 1.5  # Add line spacing
                    p.text = sentence
                    p.font.name = "Trebuchet MS"
                    p.font.size = Inches(0.25)
                    p.alignment = PP_ALIGN.CENTER
            
            # Add page number at bottom center
            page_number = page_num + 1  # Add 1 since page_num is 0-based
            page_num_box = slide.shapes.add_textbox(
                Inches(0), Inches(6.5),  # Y position near bottom of slide
                Inches(10), Inches(0.5)  # Full width of slide for centering
            )
            page_num_frame = page_num_box.text_frame
            page_num_frame.text = str(page_number)
            page_num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            page_num_frame.paragraphs[0].font.name = "Trebuchet MS"
            page_num_frame.paragraphs[0].font.size = Inches(0.15)  # Slightly smaller than main text
        
        # Save the presentation
        prs.save(str(output_path))
        print(f"Book exported to: {output_path.absolute()}")
    
    def __len__(self) -> int:
        """Return the number of pages in the book."""
        return len(self.pages)

    def illustrate(self, save_dir: str | Path | None = None) -> None:
        """
        Generate illustrations for all pages using the Hugging Face Inference API.
        
        Args:
            save_dir: Optional directory to save the generated images. 
                     If None, creates a temporary directory.
        """
        token = huggingface_hub.get_token()
        if not token:
            warnings.warn("No Hugging Face token found. Please login using `huggingface-cli login` or set the HF_TOKEN environment variable. Otherwise, you may be rate limited.")

        API_URL = f"https://api-inference.huggingface.co/models/{self.lora}"
        headers = {"Authorization": f"Bearer {token}"}

        # Create save directory if provided
        if save_dir:
            save_dir = Path(save_dir)
            save_dir.mkdir(parents=True, exist_ok=True)
        else:
            save_dir = Path(tempfile.mkdtemp())

        print("Generating illustrations... This could take a few minutes.")

        # Create a list of tasks for the progress bar
        tasks = []
        if self.title_illustration is None:
            tasks.append(("title", self.title, None))
        tasks.extend((f"page_{i+1}", text, current_illust) 
                     for i, (text, current_illust) in enumerate(zip(self.pages, self.illustrations)))

        for task_name, text, current_illust in tqdm(tasks, desc="Generating illustrations"):
            # Skip if illustration already exists or is explicitly disabled
            if isinstance(current_illust, str) or current_illust is False:
                continue
                
            try:
                print(f"\n=== Processing {task_name} ===")
                print(f"Original text: {text}")
                
                # Use custom prompt if available, otherwise get from Qwen
                if task_name == "title" and self.title_illustration_prompt:
                    illustration_prompt = self.title_illustration_prompt
                elif task_name != "title":
                    page_num = int(task_name.split('_')[1]) - 1
                    if self.illustration_prompts[page_num]:
                        illustration_prompt = self.illustration_prompts[page_num]
                    else:
                        illustration_prompt = self._get_illustration_prompt(text)
                else:
                    illustration_prompt = self._get_illustration_prompt(text)
                    
                print(f"Illustration prompt: {illustration_prompt}")
                
                # Then get the final prompt and query the image API
                prompt = self._get_prompt(illustration_prompt)
                print(f"Final image prompt: {prompt}")
                
                response = requests.post(API_URL, headers=headers, json={"inputs": prompt})
                
                if response.status_code != 200:
                    print(f"Warning: Failed to generate illustration for {task_name}: {response.text}")
                    continue
                    
                # Save the image
                image = Image.open(io.BytesIO(response.content))
                image_path = save_dir / f"{task_name}.png"
                image.save(image_path)
                print(f"Image saved to: {image_path}")
                
                # Update the appropriate illustration reference
                if task_name == "title":
                    self.title_illustration = str(image_path)
                else:
                    page_num = int(task_name.split('_')[1]) - 1
                    self.illustrations[page_num] = str(image_path)
                    
            except Exception as e:
                print(f"Warning: Error generating illustration for {task_name}: {e}")
                continue

        print(f"\nAll illustrations saved to: {save_dir}")

